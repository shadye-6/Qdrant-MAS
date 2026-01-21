import os
import time
def prof(label):
    print(f"[{time.strftime('%H:%M:%S')}] {label}", flush=True)
import fitz
from pptx import Presentation
import subprocess
from llama_index.core import Document
from utils import (
    describe_image, is_graph, process_graph, extract_text_around_item,
    process_text_blocks, save_uploaded_file
)
from agents.graph_extractor import GraphExtractor
from agents.graph_store import GraphStore

def base_metadata(source, modality, doc_type, page_num=None):
    return {
        "source": source,
        "modality": modality,
        "type": doc_type,
        "page_num": page_num,
        "timestamp": time.time()
    }

# ---------------- PDF ----------------

def get_pdf_documents(pdf_file, llm):
    prof("Opening PDF")
    all_pdf_documents = []

    try:
        f = fitz.open(stream=pdf_file.read(), filetype="pdf")
    except Exception as e:
        print(f"PDF open error: {e}")
        return []

    prof(f"PDF opened. Pages = {len(f)}")

    for i in range(len(f)):
        prof(f"Page {i+1} start")

        page = f[i]

        prof("Extracting text blocks")
        text_blocks = [
            block for block in page.get_text("blocks", sort=True)
            if block[-1] == 0
        ]

        prof(f"Text blocks: {len(text_blocks)}")

        prof("Grouping text blocks")
        grouped_text_blocks = process_text_blocks(text_blocks)
        prof(f"Grouped blocks: {len(grouped_text_blocks)}")

        # ---------- Tables ----------
        prof("Parsing tables")
        try:
            table_docs, table_bboxes, _ = parse_all_tables(
                pdf_file.name, page, i, text_blocks, llm
            )
        except Exception as e:
            print("Table parsing failed:", e)
            table_docs, table_bboxes = [], []

        prof(f"Tables extracted: {len(table_docs)}")
        all_pdf_documents.extend(table_docs)

        # ---------- Images ----------
        prof("Parsing images")
        try:
            image_docs = parse_all_images(
                pdf_file.name, page, i, text_blocks, llm
            )
        except Exception as e:
            print("Image parsing failed:", e)
            image_docs = []

        prof(f"Images extracted: {len(image_docs)}")
        all_pdf_documents.extend(image_docs)

        # ---------- Text ----------
        prof("Processing text groups")

        for ctr, (heading_block, content) in enumerate(grouped_text_blocks, 1):
            bbox = fitz.Rect(heading_block[:4])

            if not any(bbox.intersects(tb) for tb in table_bboxes):
                meta = base_metadata(
                    source=f"{pdf_file.name}-page{i}-block{ctr}",
                    modality="text",
                    doc_type="pdf_text",
                    page_num=i
                )

                all_pdf_documents.append(
                    Document(
                        text=f"{heading_block[4]}\n{content}",
                        metadata=meta
                    )
                )

        prof(f"Page {i+1} finished")

    f.close()
    prof("PDF fully processed")

    return all_pdf_documents

# ---------------- Tables ----------------
#     table_docs = []
#     table_bboxes = []

#     try:
        # tables = page.find_tables(horizontal_strategy="lines_strict", vertical_strategy="lines_strict")
        # prof("Table detection end")

#         for idx, tab in enumerate(tables):
#             if tab.header.external:
#                 continue

#             df = tab.to_pandas()

#             out_dir = os.path.join(os.getcwd(), "vectorstore/table_references")
#             os.makedirs(out_dir, exist_ok=True)

#             safe = os.path.basename(filename).replace(".", "_")
#             df_path = os.path.join(out_dir, f"table-{safe}-{pagenum}-{idx}.xlsx")
#             df.to_excel(df_path)

#             bbox = fitz.Rect(tab.bbox)
#             table_bboxes.append(bbox)

#             before_text, after_text = extract_text_around_item(text_blocks, bbox, page.rect.height)

#             table_img = page.get_pixmap(clip=bbox)
#             img_path = os.path.join(out_dir, f"table-{safe}-{pagenum}-{idx}.jpg")
#             table_img.save(img_path)

#             caption = f"{before_text} {after_text}".strip()

#             meta = base_metadata(
#                 source=f"{filename}-page{pagenum}-table{idx}",
#                 modality="table",
#                 doc_type="pdf_table",
#                 page_num=pagenum
#             )

#             meta.update({
#                 "dataframe_path": df_path,
#                 "image_path": img_path,
#                 "caption": caption
#             })

#             table_text = f"Table caption: {caption}\nColumns: {', '.join(df.columns.astype(str))}"
#             doc = Document(
#                 text=table_text,
#                 metadata=meta
#             )
#             table_docs.append(doc)
            
#             # Extract and store graph data for table
#             try:
#                 graph_data = graph_extractor.extract(table_text)
#                 graph_store.store(
#                     graph_data["entities"],
#                     graph_data["relations"],
#                     chunk_id=meta["source"]
#                 )
#             except Exception as e:
#                 print(f"Graph extraction error for table {meta['source']}: {e}")

#     except Exception as e:
#         print(f"Table extraction error: {e}")

#     return table_docs, table_bboxes, {}

def parse_all_tables(filename, page, pagenum, text_blocks, llm, graph_extractor=None, graph_store=None):
    prof("Table detection start")
    tables = page.find_tables(horizontal_strategy="lines_strict", vertical_strategy="lines_strict")
    prof("Table detection end")
    
    # ...existing code...
    return [], [], {}


# ---------------- Images ----------------

# def parse_all_images(filename, page, pagenum, text_blocks, llm, graph_extractor, graph_store):
#     image_docs = []

#     for image_info in page.get_image_info(xrefs=True):
#         xref = image_info["xref"]
#         if xref == 0:
#             continue

#         extracted = page.parent.extract_image(xref)
#         image_data = extracted["image"]

#         out_dir = os.path.join(os.getcwd(), "vectorstore/image_references")
#         os.makedirs(out_dir, exist_ok=True)

#         img_path = os.path.join(out_dir, f"{filename}-page{pagenum}-img{xref}.png")

#         with open(img_path, "wb") as f:
#             f.write(image_data)

#         caption = describe_image(image_data)

#         meta = base_metadata(
#             source=f"{filename}-page{pagenum}-image{xref}",
#             modality="image",
#             doc_type="pdf_image",
#             page_num=pagenum
#         )

#         meta.update({
#             "image_path": img_path,
#             "caption": caption
#         })

#         image_text = f"Image: {caption}"
#         doc = Document(text=image_text, metadata=meta)
#         image_docs.append(doc)
        
#         # Extract and store graph data for image
#         try:
#             graph_data = graph_extractor.extract(image_text)
#             graph_store.store(
#                 graph_data["entities"],
#                 graph_data["relations"],
#                 chunk_id=meta["source"]
#             )
#         except Exception as e:
#             print(f"Graph extraction error for image {meta['source']}: {e}")

#     return image_docs

# ---------------- PPT ----------------

def process_ppt_file(ppt_path, graph_extractor, graph_store):
    pdf_path = convert_ppt_to_pdf(ppt_path)
    images_data = convert_pdf_to_images(pdf_path)
    slide_texts = extract_text_and_notes_from_ppt(ppt_path)

    docs = []

    for (img_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
        text = slide_text + (f"\nNotes: {notes}" if notes else "")

        meta = base_metadata(
            source=os.path.basename(ppt_path),
            modality="image",
            doc_type="ppt_slide",
            page_num=page_num
        )

        meta["image_path"] = img_path

        doc = Document(text=text, metadata=meta)
        docs.append(doc)
        
        # Extract and store graph data for slide
        try:
            graph_data = graph_extractor.extract(text)
            graph_store.store(
                graph_data["entities"],
                graph_data["relations"],
                chunk_id=meta["source"]
            )
        except Exception as e:
            print(f"Graph extraction error for slide {meta['source']}: {e}")

    return docs

# ---------------- Utilities ----------------

def convert_ppt_to_pdf(ppt_path):
    out_dir = os.path.abspath("vectorstore/ppt_references")
    os.makedirs(out_dir, exist_ok=True)

    pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")

    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, ppt_path],
        check=True
    )

    return pdf_path

def convert_pdf_to_images(pdf_path):
    doc = fitz.open(pdf_path)
    out_dir = os.path.join(os.getcwd(), "vectorstore/ppt_references")
    os.makedirs(out_dir, exist_ok=True)

    images = []
    for i in range(len(doc)):
        pix = doc.load_page(i).get_pixmap()
        img_path = os.path.join(out_dir, f"{os.path.basename(pdf_path)}_{i}.png")
        pix.save(img_path)
        images.append((img_path, i))

    doc.close()
    return images

def extract_text_and_notes_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    result = []

    for slide in prs.slides:
        slide_text = " ".join([s.text for s in slide.shapes if hasattr(s, "text")])
        notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
        result.append((slide_text, notes))

    return result

# ---------------- Loaders ----------------

def load_multimodal_data(files, llm, graph_extractor, graph_store):
    documents = []

    for file in files:
        ext = os.path.splitext(file.name.lower())[1]

        if ext in [".png", ".jpg", ".jpeg"]:
            content = file.read()
            caption = describe_image(content)

            meta = base_metadata(file.name, "image", "uploaded_image")
            documents.append(Document(text=caption, metadata=meta))

        elif ext == ".pdf":
            documents.extend(get_pdf_documents(file, llm))

        elif ext in [".ppt", ".pptx"]:
            documents.extend(process_ppt_file(save_uploaded_file(file), graph_extractor, graph_store))

        else:
            text = file.read().decode("utf-8", errors="ignore")
            meta = base_metadata(file.name, "text", "text_file")
            doc = Document(text=text, metadata=meta)
            documents.append(doc)
            
            # Extract and store graph data for text file
            try:
                graph_data = graph_extractor.extract(text)
                graph_store.store(
                    graph_data["entities"],
                    graph_data["relations"],
                    chunk_id=meta["source"]
                )
            except Exception as e:
                print(f"Graph extraction error for file {meta['source']}: {e}")

    return documents
